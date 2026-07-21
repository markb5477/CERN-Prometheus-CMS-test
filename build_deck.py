#!/usr/bin/env python3
"""Build the Tuesday meeting slides from the test results into presentation.pptx."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x1d, 0x2d, 0x44)
TEAL = RGBColor(0x2a, 0x9d, 0x8f)
RED = RGBColor(0xc1, 0x12, 0x1f)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xf2, 0xf4, 0xf7)
WHITE = RGBColor(0xff, 0xff, 0xff)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def style(run, size, color=NAVY, bold=False, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Calibri"


def bar(s, color=TEAL, h=Inches(0.16)):
    r = s.shapes.add_shape(1, 0, 0, SW, h)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    return r


def header(s, title, kicker=None):
    bar(s)
    _, tf = box(s, Inches(0.6), Inches(0.5), Inches(12.1), Inches(1.0))
    p = tf.paragraphs[0]
    style(p.add_run(), 32, NAVY, bold=True)
    p.runs[0].text = title
    if kicker:
        _, kf = box(s, Inches(0.62), Inches(0.18), Inches(12), Inches(0.4))
        kp = kf.paragraphs[0]
        r = kp.add_run(); r.text = kicker.upper(); style(r, 12, TEAL, bold=True)
    line = s.shapes.add_shape(1, Inches(0.62), Inches(1.42), Inches(2.2), Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = TEAL; line.line.fill.background()


def bullets(s, items, left=Inches(0.7), top=Inches(1.75),
            width=Inches(12), size=20, gap=10):
    _, tf = box(s, left, top, width, SH - top - Inches(0.5))
    tf.word_wrap = True
    first = True
    for it in items:
        lvl = 0
        text = it
        if isinstance(it, tuple):
            lvl, text = it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.level = lvl
        if lvl == 0:
            r = p.add_run(); r.text = "▪  "; style(r, size, TEAL, bold=True)
        else:
            p.space_before = Pt(2)
            r = p.add_run(); r.text = "      –  "; style(r, size - 3, GREY)
        # allow **bold** segments
        for i, seg in enumerate(text.split("**")):
            if seg == "":
                continue
            r = p.add_run(); r.text = seg
            style(r, size if lvl == 0 else size - 3,
                  NAVY if lvl == 0 else GREY, bold=(i % 2 == 1))
    return tf


def footer(s, n=None):
    n = len(prs.slides._sldIdLst)  # this slide's 1-based position; survives insertions
    _, tf = box(s, Inches(0.5), SH - Inches(0.45), Inches(11), Inches(0.35))
    r = tf.paragraphs[0].add_run()
    r.text = "CMS Tracker monitoring  |  single Prometheus node at 1 Hz"
    style(r, 10, GREY)
    _, nf = box(s, SW - Inches(1.0), SH - Inches(0.45), Inches(0.6), Inches(0.35))
    nf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    r = nf.paragraphs[0].add_run(); r.text = str(n); style(r, 10, GREY)


def picture(s, path, left, top, height=None, width=None):
    if os.path.exists(path):
        return s.shapes.add_picture(path, left, top, height=height, width=width)
    return None


def methodnote(s, lines, left, top, width, size=11):
    rule = s.shapes.add_shape(1, left, top, Inches(1.3), Pt(2.5))
    rule.fill.solid(); rule.fill.fore_color.rgb = TEAL; rule.line.fill.background()
    _, tf = box(s, left, top + Inches(0.07), width, Inches(1.4))
    tf.word_wrap = True
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(3)
        for i, seg in enumerate(ln.split("**")):
            if seg == "":
                continue
            r = p.add_run(); r.text = seg
            style(r, size, GREY, bold=(i % 2 == 1))


# ---------------------------------------------------------------- 1  title
s = slide()
rect = s.shapes.add_shape(1, 0, 0, SW, SH)
rect.fill.solid(); rect.fill.fore_color.rgb = NAVY; rect.line.fill.background()
bar(s, TEAL, Inches(0.22))
_, tf = box(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(2.4))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Monitoring the Tracker with Prometheus"
style(r, 44, WHITE, bold=True)
p2 = tf.add_paragraph(); p2.space_before = Pt(14)
r = p2.add_run()
r.text = "Can a single node read every parameter, once a second?"
style(r, 24, RGBColor(0x9e, 0xd8, 0xd1))
_, mf = box(s, Inches(0.92), Inches(5.6), Inches(11), Inches(1.2))
for i, line in enumerate([
        "Load tests on the HPC  |  Prometheus + Avalanche",
        "Mark Brandt  |  DAQ software meeting, building 598"]):
    p = mf.paragraphs[0] if i == 0 else mf.add_paragraph()
    r = p.add_run(); r.text = line
    style(r, 16, RGBColor(0xc9, 0xd2, 0xdd), bold=(i == 0))

# ---------------------------------------------------------------- 2  question
s = slide(); header(s, "The question", "why we ran this")
bullets(s, [
    "The upgraded Tracker will expose on the order of **900k monitoring parameters** (voltages, currents, temperatures, rates).",
    "We want a fresh reading of **every one of them, once per second** (1 Hz).",
    "Prometheus is the candidate tool. The question for this study:",
    (1, "**Can one Prometheus node keep up with that load, and if not, what breaks first, the hardware or the design?**"),
    "Approach: reproduce the load on the HPC and measure, rather than guess from spec sheets.",
], top=Inches(1.9), size=21, gap=14)
footer(s, 2)

# ---------------------------------------------------------------- 3  push vs pull
s = slide(); header(s, "How Prometheus works: pull, not push", "1 minute of background")
bullets(s, [
    "**Push** (many systems): each device sends its numbers to a central server whenever it wants.",
    "**Pull** (Prometheus): the server reaches out and **scrapes** each source over HTTP on a fixed schedule.",
    (1, "the server sets the cadence, here **once per second**"),
    (1, "it instantly knows if a source is unreachable (an **up** flag per source)"),
    (1, "sources stay simple: they just expose current values, they do not need the server's address"),
    "One number sampled over time (say a temperature) is one **time series**. 1 parameter at 1 Hz = **1 value every second**.",
], top=Inches(1.9), size=20, gap=12)
footer(s, 3)

# ---------------------------------------------------------------- 4  load test
s = slide(); header(s, "What a load test is", "the method")
bullets(s, [
    "We do not have the full detector yet, so we **synthesise** a load that looks identical to Prometheus.",
    "**Avalanche** (a generator built by the Prometheus community for exactly this) stands in for the readout.",
    (1, "each fake **module** = one source Prometheus scrapes"),
    (1, "each module exposes a chosen number of **parameters**, refreshed every second"),
    "We turn the two knobs that matter, **how many modules** and **how many parameters each**, and watch what happens.",
    "**Pass/fail is simple:** every scrape finishes inside the 1 second budget, and every module stays up.",
], top=Inches(1.9), size=20, gap=11)
footer(s, 4)

# ---------------------------------------------------------------- 5  test catalog
s = slide(); header(s, "The tests we ran", "six load patterns, plus the CMS model")
tests = [
    ("Test", "What it varies", "What it tells us"),
    ("Module ramp", "25 modules, total up to 1.2M", "does the real operating load fit inside 1 s?"),
    ("Ramp", "80 thin modules, total up to 2M", "how scrape time grows as the load rises"),
    ("Sweep", "2M fixed, split across 1 to 160 modules", "the real limit is parameters per module, not total"),
    ("Stress", "80 modules, large jumps up to 2M", "find the breaking point quickly"),
    ("Spike", "baseline, burst to 2M, back to baseline", "burst survival and recovery"),
    ("Soak", "1M held for 30 minutes", "memory creep or scrape drift over time"),
    ("CMS growth", "real 316-module model, 880k to 2.5M", "headroom before the node's total ceiling"),
    ("CMS granularity", "880k fixed, 316 down to 8 modules", "how coarser modules fatten each scrape"),
]
gt = s.shapes.add_table(len(tests), 3, Inches(0.7), Inches(1.72),
                        Inches(11.95), Inches(4.9)).table
gt.first_row = False; gt.horz_banding = False
gt.columns[0].width = Inches(2.25)
gt.columns[1].width = Inches(4.35)
gt.columns[2].width = Inches(5.35)
for ri, row in enumerate(tests):
    for ci, val in enumerate(row):
        cell = gt.cell(ri, ci)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.11); cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = val
        if ri == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            style(run, 15, WHITE, bold=True)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if ri % 2 else WHITE
            style(run, 13, NAVY, bold=(ci == 0))
footer(s)

# ---------------------------------------------------------------- 6  setup
s = slide(); header(s, "The setup", "what we ran, and where")
bullets(s, [
    "**Server:** one Prometheus instance on the HPC (32 cores, 125 GB RAM).",
    "**Load:** many Avalanche exporters, one per module, scraped at 1 Hz (timeout 0.9 s).",
    "**Model:** parameter counts taken from the monitoring design note, ~880k across 316 per-DTC aggregation points.",
    "**Everything is scripted and repeatable** (bash driving both sides); the full suite runs in about an hour.",
], top=Inches(1.9), size=20, gap=13)
_, rf = box(s, Inches(0.7), Inches(5.35), Inches(12), Inches(1.0))
rp = rf.paragraphs[0]
r = rp.add_run(); r.text = "Repository:  "; style(r, 18, NAVY, bold=True)
r = rp.add_run(); r.text = "github.com/markb5477/CERN-Prometheus-CMS-test"
style(r, 18, TEAL, bold=True)
footer(s, 5)

# ---------------------------------------------------------------- 6  headline result
s = slide(); header(s, "Headline: the real load is easy", "finding 0")
bullets(s, [
    "Real model: **880k parameters across 316 modules** (about 2,800 each).",
    "Result: every scrape completes in **0.03 s**, all 316 modules up.",
    (1, "that is roughly **3% of the one-second budget**"),
    "So at face value, yes, a single node handles the Tracker's real monitoring load comfortably.",
    "The interesting part is **how much margin there is, and how it fails** when pushed. Two limits showed up.",
], top=Inches(1.9), size=21, gap=14)
footer(s, 6)

# ---------------------------------------------------------------- 7  limit 1
s = slide(); header(s, "Limit 1: parameters per module, not the total", "finding 1")
bullets(s, [
    "The same total load behaves very differently depending on how it is **spread**.",
    "Take **2 million parameters** and split them different ways:",
    (1, "over **40 modules** (50k each): **0.81 s**, all up  ✓"),
    (1, "over **20 modules** (100k each): **1.10 s**, only **7 of 20** up  ✗"),
    (1, "in **one 2M module**: never even comes up  ✗"),
    "It is not the volume that hurts, it is **concentration**. Design rule: keep each scrape target modest.",
], left=Inches(0.7), top=Inches(1.85), width=Inches(7.2), size=18, gap=9)
picture(s, "results/suite.png", Inches(8.1), Inches(1.7), Inches(5.3))
methodnote(s, [
    "**How to read:** each point is a **single reading, 45 s after the load starts**, "
    "plotting the **slowest module's** scrape. Not averaged, not repeated (N=1); "
    "run-to-run jitter ~5% (from the soak). Dashed line = 1 s budget.",
], Inches(0.7), Inches(5.95), Inches(7.2))
footer(s, 7)

# ---------------------------------------------------------------- 8  limit 2
s = slide(); header(s, "Limit 2: the node's total ceiling", "finding 2")
bullets(s, [
    "Hold the real 316-module layout, push the **total** parameter count up:",
    (1, "880k (real): 0.03 s  ✓"),
    (1, "2.0M: 0.27 s, all 316 up  ✓"),
    (1, "2.5M: node **collapses**, only 15 of 316 up  ✗"),
    "The single-node ceiling sits **between 2M and 2.5M**.",
    "Against the 880k real load, that is about **2.3x headroom**.",
], left=Inches(0.7), top=Inches(1.95), width=Inches(5.7), size=18, gap=10)
picture(s, "results/cms.png", Inches(6.7), Inches(2.55), width=Inches(6.4))
methodnote(s, [
    "**How to read:** each point is a **single reading, 45 s after the load starts** "
    "(worst-case scrape across modules); N=1, run-to-run jitter ~5%.",
    "**The 2.5M dot looks fast but only 15/316 are up:** most targets are down, so the "
    "time reflects only the survivors, not health.",
], Inches(0.7), Inches(5.75), Inches(12.2))
footer(s, 8)

# ---------------------------------------------------------------- 9  stability
s = slide(); header(s, "It is also stable, not just fast once", "finding 3")
bullets(s, [
    "**Sustained:** 1M parameters held for **30 minutes**. Scrape time stayed flat at ~0.17 s, no drift, all modules up throughout.",
    "**Bursts:** jump 400k → 2M → back to 400k. The node absorbs the spike (0.48 s) and **recovers cleanly**.",
    "**Granularity check:** at fixed 880k, merging into fewer/fatter modules (down to 8 modules, 110k each) pushes scrape to **0.93 s**, right at the budget, which reinforces Limit 1.",
    (1, "one honest caveat: memory crept up over the 30-minute hold (trivial on 125 GB, but a longer soak at real load is on the follow-up list)."),
], top=Inches(1.9), size=19, gap=13)
footer(s, 9)

# ---------------------------------------------------------------- 10  scope
s = slide(); header(s, "What we tested, and what we did not", "scope, being honest")
_, lh = box(s, Inches(0.7), Inches(1.7), Inches(5.6), Inches(0.4))
r = lh.paragraphs[0].add_run(); r.text = "Validated on the HPC"; style(r, 17, TEAL, bold=True)
_, rh = box(s, Inches(6.95), Inches(1.7), Inches(5.9), Inches(0.4))
r = rh.paragraphs[0].add_run(); r.text = "Not yet tested"; style(r, 17, RED, bold=True)
bullets(s, [
    "**880k** real load at 1 Hz: **0.03 s**, all up",
    "Ingestion healthy to **~2M series** (~2.3x headroom)",
    "Per-module limit found: **thin targets win**",
    "30-minute hold stable, burst recovery clean",
], left=Inches(0.7), top=Inches(2.2), width=Inches(5.6), size=16, gap=11)
bullets(s, [
    "**Series churn** (firmware renames) - the real Prometheus limit",
    "**Realistic noisy values** (compression, deadband)",
    "**Cold tier**: downsampling / long-term / remote-write",
    "**Long soak**: disk growth, days, memory creep",
    "**Query + alerting** load (ingest-only so far)",
    "**Production hardware + network** (was one box, localhost)",
    "**HA failover / sharding** (proposed, not proven)",
], left=Inches(6.95), top=Inches(2.2), width=Inches(5.9), size=16, gap=8)
footer(s, 10)

# ---------------------------------------------------------------- 11  solutions
s = slide(); header(s, "If we deploy: two ways to make it production-grade", "solutions")
bullets(s, [
    "One node already meets the real load with margin, but a single node is a **single point of failure**.",
    "**Option A, twin nodes (redundancy):** two identical Prometheus nodes scrape the same targets. If one dies, the other still has the full picture. Standard Prometheus HA, no data gaps.",
    "**Option B, functional sharding:** split the modules across several smaller nodes, each owning a fixed range, each with a twin. Scales past the single-node ceiling and limits the blast radius if one fails.",
    "For the real load this is about **reliability, not capacity**, one node has the throughput. Both options still need a **failover test**.",
], top=Inches(1.9), size=19, gap=14)
footer(s, 11)

# ---------------------------------------------------------------- 12  storage open q
def dbox(s, left, top, w, h, lines, border=TEAL, fill=WHITE, tsize=11.5):
    shp = s.shapes.add_shape(5, left, top, w, h)      # rounded rectangle
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border; shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        for j, seg in enumerate(ln.split("**")):
            if seg == "":
                continue
            r = p.add_run(); r.text = seg; style(r, tsize, NAVY, bold=(j % 2 == 1))
    return shp

def flowlabel(s, left, top, w, text, color=GREY):
    _, tf = box(s, left, top, w, Inches(0.3))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; style(r, 11, color, bold=True)

s = slide(); header(s, "Open question: where the data lives, and for how long",
                     "storage and retention, left for integration")
bullets(s, [
    "**It always lives on disk.** Prometheus writes compressed blocks to local disk (~1-2 bytes/sample), not RAM. Nothing is silently dropped.",
    "**It does delete, by design.** Set a retention window (say 15 days); past that, Prometheus removes old data. That is why the hot tier is 'short retention'.",
    "**What our tests did (for simplicity):** wrote to local disk and **wiped it between runs**. Retention and any cold tier were **not** exercised.",
    "**So the open question is the handoff:** what to keep, and where, before Prometheus would otherwise delete it.",
], left=Inches(0.6), top=Inches(1.8), width=Inches(6.0), size=15, gap=10)

dx, dw = Inches(6.9), Inches(6.0)
dbox(s, dx, Inches(1.8), dw, Inches(0.68),
     ["**Monitoring layer**", "1 Hz, ~1M parameters"], border=GREY)
flowlabel(s, dx, Inches(2.54), dw, "↓   ~30 Mb/s firehose")
dbox(s, dx, Inches(2.9), dw, Inches(1.0),
     ["**HOT TIER  -  Prometheus / VictoriaMetrics**",
      "local disk, full 1 Hz, short window (days to weeks)",
      "auto-deletes past retention"])
flowlabel(s, dx, Inches(3.96), dw,
          "↓   reduce here (OPEN): deadband / downsample / keep-all", color=RED)
dbox(s, dx, Inches(4.4), dw, Inches(1.25),
     ["**COLD TIER  -  long-term archive**",
      "ASIC params  ->  relational DB (joinable w/ config & calib)",
      "non-ASIC  ->  Prometheus-compatible (schema-flexible)"], border=NAVY)

methodnote(s, [
    "**How long, and at what granularity after reduction, is still open** - a physics / "
    "operations policy (how far back to debug or correlate with a run), not a DB decision. "
    "The tiering just makes any policy affordable; real thresholds come from rates measured "
    "at integration.",
], Inches(0.6), Inches(6.0), Inches(12.3))
footer(s)

# ---------------------------------------------------------------- 13  recommendation
s = slide(); header(s, "Where this leaves us", "recommendation and next steps")
bullets(s, [
    "**Feasible so far:** one node reads the real load at 1 Hz using ~3% of the budget, ~2.3x headroom, on lab hardware.",
    "**Guardrails from the tests:** keep parameters-per-module modest, and stay under ~2M series per node.",
    "**Before committing, close the gaps that raw ingestion volume hides:**",
    (1, "long soak at 880k over hours/days: memory, disk, compaction"),
    (1, "series churn from firmware updates, not flat gauges"),
    (1, "realistic values plus query/alerting load, not ingest alone"),
    (1, "re-measure on the production node, over the network"),
    (1, "cold tier, and a VictoriaMetrics head-to-head if churn bites"),
], top=Inches(1.85), size=19, gap=9)
footer(s, 12)

out = "presentation.pptx"
prs.save(out)
print("->", out, f"({len(prs.slides._sldIdLst)} slides)")
